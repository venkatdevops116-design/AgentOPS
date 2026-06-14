"""
AgentOPS Technical Presentation Generator - Scratch Level
Comprehensive technical presentation covering:
- What is AIOps (basics for zero knowledge)
- What is AgenticOps (basics for zero knowledge)
- Key differences between AIOps and AgenticOps
- Why we use both
- Real POC demonstration
- Demo scenarios
- Technical deep dive
- Architecture explained simply
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_technical_presentation():
    """Create comprehensive technical presentation for zero-knowledge audience"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    DARK_BLUE = RGBColor(31, 78, 121)
    ACCENT_BLUE = RGBColor(0, 102, 204)
    ACCENT_GREEN = RGBColor(0, 176, 80)
    ACCENT_ORANGE = RGBColor(255, 102, 0)
    ACCENT_RED = RGBColor(192, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(242, 242, 242)
    DARK_GRAY = RGBColor(89, 89, 89)
    
    def add_title_slide(prs, title, subtitle=""):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            for line in subtitle.split('\n'):
                p = subtitle_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, title, content_items):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE
        
        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        return slide
    
    def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
        """Add two-column comparison slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE
        
        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        
        # Left column header
        left_header = slide.shapes.add_textbox(Inches(0.3), Inches(1), Inches(4.5), Inches(0.4))
        left_header_frame = left_header.text_frame
        p = left_header_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        # Left content
        left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, item in enumerate(left_items):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)
            p.space_after = Pt(4)
        
        # Right column header
        right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1), Inches(4.5), Inches(0.4))
        right_header_frame = right_header.text_frame
        p = right_header_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        # Right content
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(right_items):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)
            p.space_after = Pt(4)
        
        return slide
    
    # ========== SLIDE 1: TITLE ==========
    add_title_slide(prs, "AIOps vs AgenticOps", 
                   "Understanding Modern Intelligent Operations\nFrom Zero to Advanced Concepts\nTechnical Deep Dive")
    
    # ========== SLIDE 2: WHAT IS OPERATIONS? (BASICS) ==========
    add_content_slide(prs, "What is Operations? (The Foundation)", [
        "📊 Operations = Keeping systems running smoothly 24/7",
        "",
        "Traditional Operations Work:",
        "   1️⃣  Monitor: Watch servers, databases, applications",
        "   2️⃣  Alert: When something goes wrong, sound alarm",
        "   3️⃣  Respond: Engineer investigates the problem",
        "   4️⃣  Fix: Apply solution manually",
        "   5️⃣  Verify: Check if system recovered",
        "",
        "Example Real Scenario:",
        "   • Website suddenly gets slow at 3 PM",
        "   • Alert fires: 'High CPU Usage'",
        "   • On-call engineer wakes up, checks logs (20 min)",
        "   • Finds: 'Database query too slow'",
        "   • Fixes: Optimizes query or adds cache",
        "   • System recovers (45-60 minutes total)",
        "",
        "The Problem: This is SLOW, MANUAL, and REQUIRES HUMANS"
    ])
    
    # ========== SLIDE 3: THE OPERATIONS PAIN POINTS ==========
    add_content_slide(prs, "Why Traditional Operations is Hard", [
        "😫 Current Pain Points:",
        "",
        "⏱️  SLOW RESPONSE:",
        "   • Takes 45-60 minutes to fix common issues",
        "   • Each minute of downtime = lost money/customers",
        "",
        "👤 HUMAN DEPENDENT:",
        "   • Needs skilled engineers on-call 24/7",
        "   • Engineers get tired, make mistakes",
        "   • Expensive to hire and retain experts",
        "",
        "🔴 REACTIVE (Not Proactive):",
        "   • We wait for problem to happen, then react",
        "   • No prediction or prevention",
        "",
        "😤 BURNOUT:",
        "   • On-call duty every night is exhausting",
        "   • High stress = poor decisions",
        "   • Engineers quit from burnout",
        "",
        "❌ INCONSISTENT:",
        "   • Different engineers handle issues differently",
        "   • No standardized solutions"
    ])
    
    # ========== SLIDE 4: WHAT IS AIOPS? (SCRATCH LEVEL) ==========
    add_content_slide(prs, "What is AIOps? (Simple Explanation)", [
        "AIOps = Artificial Intelligence + Operations",
        "Pronounced: \"Ay-Ops\" or \"A-I-Ops\"",
        "",
        "🤖 Think of it like this:",
        "   • You have a super-smart assistant watching your systems",
        "   • This assistant ANALYZES and DIAGNOSES problems",
        "   • But it ONLY SUGGESTS solutions, doesn't execute",
        "",
        "Real Example:",
        "   Problem: Website is slow",
        "   Traditional Op: Engineer manually checks logs (1 hour)",
        "   AIOps: AI analyzes logs instantly (2 minutes) and says:",
        "      'I found the problem! It's DATABASE SLOW.",
        "      I recommend: Optimize query OR add cache.'",
        "",
        "🎯 KEY POINT: AIOps is about ANALYSIS & RECOMMENDATIONS",
        "   • It DIAGNOSES the problem",
        "   • It SUGGESTS the solution",
        "   • HUMANS must approve and execute",
        "",
        "⏳ Speed improvement: 60 min → 10-15 min (for analysis)"
    ])
    
    # ========== SLIDE 5: WHAT IS AGENTICOPS? (SCRATCH LEVEL) ==========
    add_content_slide(prs, "What is AgenticOps? (Simple Explanation)", [
        "AgenticOps = Agentic AI + Operations",
        "Pronounced: \"Uh-JEN-tick Ops\"",
        "",
        "🤖 What's different from AIOps?",
        "   • AIOps = analyzes and RECOMMENDS",
        "   • AgenticOps = analyzes and ACTS (executes automatically)",
        "",
        "Think of it like this:",
        "   • AIOps = GPS app that says 'Turn left here'",
        "   • AgenticOps = Self-driving car that TURNS automatically",
        "",
        "Real Example (Same problem as before):",
        "   Problem: Website is slow",
        "   AgenticOps: AI analyzes logs instantly (2 minutes) and:",
        "      1. DECIDES: 'This is a slow database query'",
        "      2. EXECUTES: Auto-optimizes query OR adds cache",
        "      3. VERIFIES: Checks if website is fast now",
        "      4. NOTIFIES: Sends report to team",
        "   All AUTOMATIC without human approval!",
        "",
        "⚡ Speed improvement: 60 min → 2-3 min (entire cycle automated)"
    ])
    
    # ========== SLIDE 6: SIDE BY SIDE COMPARISON ==========
    add_two_column_slide(prs, "AIOps vs AgenticOps - Head to Head",
        "🔍 AIOps (Analysis & Insights)", [
            "✓ What it does:",
            "  • Collects data from systems",
            "  • Analyzes patterns",
            "  • Identifies root causes",
            "  • Suggests solutions",
            "",
            "✓ What HUMANS do:",
            "  • Reviews suggestions",
            "  • Approves actions",
            "  • Executes fixes manually",
            "",
            "✓ Speed: 10-15 min analysis",
            "",
            "✓ Control: Human-approved",
            "",
            "✓ Risk: Low (humans review)",
            "",
            "✓ Example Tool:",
            "  Splunk, DataDog, New Relic"
        ],
        "⚙️ AgenticOps (Action & Autonomy)", [
            "✓ What it does:",
            "  • Collects data from systems",
            "  • Analyzes patterns",
            "  • Makes decisions",
            "  • EXECUTES actions automatically",
            "",
            "✓ What HUMANS do:",
            "  • Set guardrails (rules)",
            "  • Monitor results",
            "  • Intervene if needed",
            "",
            "✓ Speed: 2-3 min end-to-end",
            "",
            "✓ Control: AI-autonomous",
            "",
            "✓ Risk: Medium (guardrails)",
            "",
            "✓ Example:",
            "  Our AgentOPS POC"
        ]
    )
    
    # ========== SLIDE 7: WHEN DO WE USE EACH? ==========
    add_content_slide(prs, "When to Use AIOps vs AgenticOps?", [
        "🔍 USE AIOps WHEN:",
        "   • Problem is CRITICAL (financial transactions, patient data)",
        "   • Need human oversight and approval",
        "   • Compliance rules require manual sign-off",
        "   • You're just starting out (less risk)",
        "",
        "Examples:",
        "   ✓ Bank detecting fraud (needs human review)",
        "   ✓ Healthcare system alerts (needs doctor approval)",
        "   ✓ Finance app (compliance-heavy, needs audits)",
        "",
        "⚙️ USE AgenticOps WHEN:",
        "   • Problem is ROUTINE and well-defined",
        "   • Action is reversible or low-risk",
        "   • Speed is critical (every second matters)",
        "   • You have guardrails and monitoring",
        "",
        "Examples:",
        "   ✓ Auto-scaling servers (easily reversible)",
        "   ✓ Cache refresh (no data loss)",
        "   ✓ Log rotation (routine maintenance)",
        "   ✓ Restart service (standard runbook)",
        "",
        "🎯 BEST PRACTICE: Use BOTH together!"
    ])
    
    # ========== SLIDE 8: WHY USE BOTH? ==========
    add_content_slide(prs, "Why We Use BOTH AIOps and AgenticOps", [
        "💡 The Smart Approach = AIOps + AgenticOps Together",
        "",
        "🏗️ Architecture:",
        "   • AgenticOps handles ROUTINE issues (80% of incidents)",
        "   • AIOps analyzes COMPLEX issues (20% of incidents)",
        "   • Humans supervise CRITICAL decisions",
        "",
        "Real Workflow:",
        "   1. Problem occurs → Both systems analyze it",
        "   2. Simple issue (e.g., 'slow DB'):",
        "      → AgenticOps EXECUTES fix automatically ✅",
        "   3. Complex issue (e.g., 'weird network pattern'):",
        "      → AIOps RECOMMENDS fix, human approves 👤",
        "   4. Critical issue (e.g., 'security breach'):",
        "      → Both alert, human takes control 🚨",
        "",
        "📊 Benefits of Both:",
        "   ✅ Speed: Fast for routine issues",
        "   ✅ Safety: Humans control critical decisions",
        "   ✅ Scalability: Handle more issues with fewer staff",
        "   ✅ Learning: AI learns from both types",
        "   ✅ Flexibility: Choose approach per issue"
    ])
    
    # ========== SLIDE 9: OUR POC - AGENTOPS PROJECT ==========
    add_content_slide(prs, "Our POC: AgentOPS Project Overview", [
        "🎯 What We Built:",
        "   A working prototype combining AIOps + AgenticOps",
        "",
        "📍 Location:",
        "   Repository: github.com/venkatdevops116-design/AgentOPS",
        "",
        "🔧 What It Does:",
        "   • Monitors a Flask web application",
        "   • Collects metrics (requests, latency, errors)",
        "   • AI analyzes metrics in real-time",
        "   • AI makes decisions and takes actions",
        "   • Shows everything in dashboards",
        "",
        "💻 Technology Stack:",
        "   ✓ Flask (Python web app)",
        "   ✓ Prometheus (metrics database)",
        "   ✓ Grafana (dashboards)",
        "   ✓ Ollama + Llama 3.2 (Local AI, no cloud)",
        "   ✓ Docker (easy deployment)",
        "",
        "🎓 Why this POC?",
        "   • Simple enough to understand concepts",
        "   • Real enough to see how it works",
        "   • Educational for learning AIOps/AgenticOps"
    ])
    
    # ========== SLIDE 10: POC ARCHITECTURE DIAGRAM ==========
    add_content_slide(prs, "AgentOPS Architecture - Simple Diagram", [
        "📊 SYSTEM ARCHITECTURE:",
        "",
        "   ┌─────────────────────────────────────────┐",
        "   │   STEP 1: APPLICATION LAYER             │",
        "   │   Flask Web App (Port 5000)             │",
        "   │   • Serves web requests                 │",
        "   │   • Generates metrics                   │",
        "   └────────────────┬────────────────────────┘",
        "                    │ (metrics emitted)",
        "   ┌────────────────▼────────────────────────┐",
        "   │   STEP 2: MONITORING LAYER              │",
        "   │   Prometheus (Port 9090)                │",
        "   │   • Collects metrics every 5 seconds    │",
        "   │   • Stores time-series data             │",
        "   │   Grafana (Port 3000)                   │",
        "   │   • Shows dashboards                    │",
        "   └────────────────┬────────────────────────┘",
        "                    │ (queries metrics)",
        "   ┌────────────────▼────────────────────────┐",
        "   │   STEP 3: AI INTELLIGENCE LAYER         │",
        "   │   Agent AI (Python)                     │",
        "   │   • Queries Prometheus                  │",
        "   │   • Sends data to LLM                   │",
        "   │   • Gets AI analysis                    │",
        "   └────────────────┬────────────────────────┘",
        "                    │ (AI decisions)",
        "   ┌────────────────▼────────────────────────┐",
        "   │   STEP 4: ACTION LAYER                  │",
        "   │   Auto-Actions                          │",
        "   │   • Scale servers, fix configs, etc     │",
        "   └─────────────────────────────────────────┘"
    ])
    
    # ========== SLIDE 11: DATA FLOW IN REAL TIME ==========
    add_content_slide(prs, "Real-Time Data Flow Example", [
        "⏱️ TIMELINE: When a Problem Occurs",
        "",
        "T=0s:   📱 Requests spike (1,000 requests arrive)",
        "",
        "T=5s:   📊 Prometheus scrapes metrics",
        "        • Captures: request_count = 1000, CPU = 85%",
        "",
        "T=10s:  🤖 AI Agent queries Prometheus",
        "        • Gets latest metrics from database",
        "",
        "T=11s:  🧠 AI Analysis by LLM (Local Llama 3.2)",
        "        • LLM analyzes: 'High requests + High CPU'",
        "        • LLM decides: 'Need to scale servers'",
        "",
        "T=12s:  ⚡ Agent Executes Action",
        "        • Auto-scales: 2 servers → 4 servers",
        "        • Sends alert: 'Scaled due to traffic spike'",
        "",
        "T=13s:  ✅ System Recovers",
        "        • Load distributed across 4 servers",
        "        • CPU drops to 45%",
        "        • Response time back to normal",
        "",
        "📈 Total time: 13 seconds (vs 45-60 min manual)",
        "🎯 This is AgenticOps in action!"
    ])
    
    # ========== SLIDE 12: POC DEMO SCENARIOS ==========
    add_content_slide(prs, "POC Demo Scenarios - What We Can Show", [
        "🌊 DEMO 1: Traffic Spike Detection",
        "   • Send 500 rapid requests to Flask app",
        "   • Prometheus detects spike",
        "   • Agent AI analyzes: 'High load detected'",
        "   • Auto-action: Scales from 2→4 instances",
        "   • Grafana shows before/after metrics",
        "",
        "🔴 DEMO 2: Error Rate Analysis",
        "   • Simulate application errors (5% error rate)",
        "   • Agent analyzes: 'Database connection pool exhausted'",
        "   • Auto-action: Increases pool size 50→100",
        "   • Error rate drops to 0.1%",
        "",
        "📉 DEMO 3: Performance Degradation",
        "   • Introduce slow database query",
        "   • Agent detects latency spike",
        "   • Recommends query optimization OR caching",
        "   • Shows metrics improvement in Grafana",
        "",
        "✅ DEMO 4: Health Report Generation",
        "   • Agent generates AI-powered system health report",
        "   • Includes: metrics, trends, recommendations",
        "   • All generated automatically by AI"
    ])
    
    # ========== SLIDE 13: HOW AI MAKES DECISIONS ==========
    add_content_slide(prs, "How AI Makes Intelligent Decisions", [
        "🧠 THE AI DECISION PROCESS (Simplified):",
        "",
        "Step 1: COLLECT CONTEXT",
        "   • CPU usage: 85%",
        "   • Memory: 72%",
        "   • Request rate: 1000/sec",
        "   • Error rate: 0.2%",
        "   • Response time: 2000ms (slow)",
        "",
        "Step 2: AI ANALYSIS",
        "   Agent sends to LLM:",
        "   'High CPU, high memory, slow response time,",
        "    normal error rate. What's the problem?'",
        "",
        "Step 3: LLM REASONING",
        "   Llama 3.2 thinks:",
        "   'This looks like: resource exhaustion",
        "    Most likely cause: Too many requests",
        "    Solution: Scale servers OR optimize code'",
        "",
        "Step 4: AI DECISION",
        "   LLM outputs:",
        "   'DECISION: Scale from 2→4 servers",
        "    CONFIDENCE: 92%",
        "    REASONING: CPU is bottleneck'",
        "",
        "Step 5: EXECUTION",
        "   Agent executes: docker-compose scale web=4",
        "",
        "🎯 This is Agentic AI - making actual decisions!"
    ])
    
    # ========== SLIDE 14: LOCAL AI VS CLOUD AI ==========
    add_content_slide(prs, "Why We Use Local AI (Llama 3.2)?", [
        "🔒 LOCAL AI Benefits:",
        "",
        "✅ PRIVACY & SECURITY:",
        "   • Your data stays in your data center",
        "   • No sending sensitive metrics to cloud",
        "   • GDPR and compliance friendly",
        "",
        "✅ NO CLOUD DEPENDENCY:",
        "   • Works even if internet is down",
        "   • No reliance on AWS/OpenAI API",
        "   • No API rate limits or throttling",
        "",
        "✅ COST EFFECTIVE:",
        "   • One-time download (3GB for Llama 3.2)",
        "   • No API charges per request",
        "   • Free to run after installation",
        "",
        "✅ LATENCY:",
        "   • Local: 1-2 seconds response time",
        "   • Cloud: 3-5 seconds (includes network)",
        "",
        "⚙️ HOW IT WORKS:",
        "   • Ollama = Framework for running local LLMs",
        "   • Llama 3.2 = Actual AI model (8K tokens context)",
        "   • Runs in Docker container",
        "   • Agent talks to Ollama via HTTP API",
        "",
        "🎯 This is Enterprise-Grade, Self-Hosted AI"
    ])
    
    # ========== SLIDE 15: MONITORING & GUARDRAILS ==========
    add_content_slide(prs, "How We Keep AI Safe (Guardrails)", [
        "🛡️ GUARDRAILS = Safety Rules for AI",
        "",
        "Rule #1: RESOURCE LIMITS",
        "   • Max scale: Don't go above 10 servers",
        "   • Max memory allocation: 80GB",
        "   • Prevents runaway escalation",
        "",
        "Rule #2: ACTION VALIDATION",
        "   • AI proposes action",
        "   • System validates: 'Is this safe?'",
        "   • Only execute if safe",
        "",
        "Rule #3: HUMAN REVIEW",
        "   • Critical decisions need approval",
        "   • Logs all AI decisions for audit trail",
        "   • Team can override AI decisions",
        "",
        "Rule #4: MONITORING",
        "   • Continuous watching of AI results",
        "   • If action fails, automatic rollback",
        "   • Metrics show success/failure",
        "",
        "Rule #5: RATE LIMITING",
        "   • Limit same action: max 5 times/hour",
        "   • Prevents AI from repeating mistakes",
        "",
        "✅ Result: Safe, controlled autonomous operations",
        "⚡ Balance: Speed + Safety"
    ])
    
    # ========== SLIDE 16: REAL WORLD APPLICATIONS ==========
    add_content_slide(prs, "Real-World Applications Beyond POC", [
        "💼 WHERE THIS IS USED TODAY:",
        "",
        "🏦 BANKING & FINANCE:",
        "   • Detect fraud patterns automatically",
        "   • Adapt security rules in real-time",
        "   • Auto-trigger alerts for suspicious activity",
        "",
        "🏥 HEALTHCARE:",
        "   • Monitor patient vitals systems",
        "   • Alert on anomalies (requires doctor approval)",
        "   • Optimize resource allocation",
        "",
        "🌐 E-COMMERCE & RETAIL:",
        "   • Handle traffic spikes during sales",
        "   • Auto-scale checkout systems",
        "   • Predict and prevent outages",
        "",
        "☁️ CLOUD INFRASTRUCTURE:",
        "   • Auto-healing of failed services",
        "   • Predict server failures before they happen",
        "   • Optimize cloud costs automatically",
        "",
        "🚀 SaaS PLATFORMS:",
        "   • Multi-tenant platform management",
        "   • Isolation and resource fairness",
        "   • Self-healing microservices",
        "",
        "🎯 All use same concepts as our POC!"
    ])
    
    # ========== SLIDE 17: KEY DIFFERENCES SUMMARY ==========
    add_content_slide(prs, "Key Differences - Quick Reference", [
        "Aspect          | Traditional  | AIOps        | AgenticOps",
        "─────────────────|──────────────|──────────────|──────────────",
        "Detection Time  | 30-45 min    | 5-10 min     | 2-5 min",
        "Analysis Time   | Manual/slow  | 5-10 min     | 1-2 min",
        "Fix Time        | 30-45 min    | 1-2 min      | <1 min",
        "Total MTTR      | 60 min       | 15-20 min    | 3-5 min",
        "",
        "Decision Making | Human only   | Hybrid       | Hybrid+Auto",
        "Automation      | 0%           | 30%          | 70%",
        "Human Effort    | Very High    | Medium       | Low",
        "24/7 Coverage   | Expensive    | Possible     | Recommended",
        "",
        "Scalability     | Per engineer | Team scaling | AI scaling",
        "Consistency     | Varies       | Better       | Excellent",
        "Learning Curve  | Long         | Medium       | Medium",
        "",
        "🎯 AgentOPS achieves best of both!"
    ])
    
    # ========== SLIDE 18: QUICK START GUIDE ==========
    add_content_slide(prs, "Quick Start: Running AgentOPS POC", [
        "✅ PREREQUISITES:",
        "   • Docker installed (docker --version)",
        "   • Python 3.11+ installed (python --version)",
        "   • 4GB+ available disk space",
        "",
        "📥 STEP 1: Clone Repository (1 minute)",
        "   git clone https://github.com/venkatdevops116-design/AgentOPS.git",
        "   cd AgentOPS",
        "",
        "🚀 STEP 2: Start Services (2 minutes)",
        "   docker compose up --build -d",
        "",
        "✓ STEP 3: Verify Services Running (1 minute)",
        "   • Flask: curl http://localhost:5000",
        "   • Prometheus: http://localhost:9090",
        "   • Grafana: http://localhost:3000",
        "",
        "🤖 STEP 4: Run AI Agent (1 minute)",
        "   python agentic_agent.py",
        "",
        "📊 STEP 5: Generate Traffic (1 minute)",
        "   for i in {1..500}; do curl http://localhost:5000; done",
        "",
        "⏱️ TOTAL TIME: ~5 minutes to see AgenticOps in action!"
    ])
    
    # ========== SLIDE 19: HANDS-ON DEMO WALKTHROUGH ==========
    add_content_slide(prs, "Live Demo Walkthrough", [
        "🎬 WHAT YOU'LL SEE IN DEMO:",
        "",
        "SCREEN 1: Grafana Dashboard",
        "   • Real-time graphs of metrics",
        "   • Request count, CPU usage, response time",
        "   • All updating live",
        "",
        "SCREEN 2: Terminal with AI Agent",
        "   • Agent prints: 'Analyzing system metrics...'",
        "   • Shows: 'Found issue: High traffic detected'",
        "   • Outputs: 'Decision: Scale servers 2→4'",
        "   • Action: 'Scaling... Done!'",
        "",
        "SCREEN 3: Prometheus Queries",
        "   • Live PromQL queries",
        "   • Shows raw metrics data",
        "   • What AI actually analyzes",
        "",
        "RESULT:",
        "   • Before: High latency, CPU 95%",
        "   • After: Low latency, CPU 45%",
        "   • AI made the decision autonomously",
        "",
        "🎯 This demonstrates AgenticOps in real-time!"
    ])
    
    # ========== SLIDE 20: LEARNING OUTCOMES ==========
    add_content_slide(prs, "What You've Learned Today", [
        "✅ UNDERSTANDING:",
        "   1. What Operations means and why it's hard",
        "   2. What AIOps is (analysis + recommendations)",
        "   3. What AgenticOps is (autonomous execution)",
        "   4. Key differences between both approaches",
        "   5. Why we use both together",
        "",
        "✅ TECHNICAL CONCEPTS:",
        "   6. How metrics flow through systems",
        "   7. How AI analyzes operational data",
        "   8. How decisions are made autonomously",
        "   9. Guardrails and safety mechanisms",
        "   10. Local vs Cloud AI trade-offs",
        "",
        "✅ PRACTICAL KNOWLEDGE:",
        "   11. AgentOPS POC architecture",
        "   12. Technology stack used",
        "   13. How to run the POC yourself",
        "   14. Real-world applications",
        "   15. Demo scenarios and live examples",
        "",
        "🎓 You now understand modern intelligent operations!"
    ])
    
    # ========== SLIDE 21: NEXT STEPS & QUESTIONS ==========
    add_content_slide(prs, "Next Steps & Q&A", [
        "🔗 RESOURCES:",
        "   Repository: github.com/venkatdevops116-design/AgentOPS",
        "   Documentation: README.md in repo",
        "   Technology Docs:",
        "   • Prometheus: prometheus.io",
        "   • Grafana: grafana.com",
        "   • Ollama: ollama.ai",
        "   • Flask: flask.palletsprojects.com",
        "",
        "📚 TO LEARN MORE:",
        "   1. Run the POC yourself",
        "   2. Modify the AI prompts and see different decisions",
        "   3. Add new metrics to monitor",
        "   4. Create new demo scenarios",
        "   5. Integrate with your own applications",
        "",
        "❓ QUESTIONS?",
        "   What would you like to know more about?",
        "   • How to customize AI decisions?",
        "   • How to add more metrics?",
        "   • How to integrate with existing systems?",
        "   • Real-world deployment considerations?",
        "",
        "Thank you! 🙏"
    ])
    
    return prs

def save_presentation(prs, filename='AgentOPS_Technical_Presentation.pptx'):
    """Save presentation to file"""
    prs.save(filename)
    return filename

if __name__ == '__main__':
    print("🎬 Generating AgentOPS Technical Presentation (Scratch Level)...")
    prs = create_technical_presentation()
    filename = save_presentation(prs)
    print(f"✅ Presentation saved: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\nPresentation includes:")
    print("  ✓ Scratch-level technical content")
    print("  ✓ What is AIOps (basics)")
    print("  ✓ What is AgenticOps (basics)")
    print("  ✓ Key differences explained simply")
    print("  ✓ Why we use both")
    print("  ✓ POC architecture and components")
    print("  ✓ Real-time data flow examples")
    print("  ✓ Demo scenarios")
    print("  ✓ AI decision-making process")
    print("  ✓ Safety guardrails")
    print("  ✓ Real-world applications")
    print("  ✓ Quick start guide")
    print("  ✓ Live demo walkthrough")
    print("  ✓ Learning outcomes")
    print("\n🎯 Perfect for audiences with zero prior knowledge!")
